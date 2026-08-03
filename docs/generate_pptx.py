import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    bg_dark = RGBColor(11, 15, 25)
    card_bg = RGBColor(18, 26, 43)
    accent_gold = RGBColor(251, 191, 36)
    accent_cyan = RGBColor(56, 189, 248)
    accent_green = RGBColor(16, 185, 129)
    text_light = RGBColor(248, 250, 252)
    text_muted = RGBColor(148, 163, 184)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_dark

    def add_header(slide, title_text, category_text="HACKATHON ACADÉMIQUE IPM 2026 - GROUPE 3"):
        # Header category
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = accent_cyan

        # Header Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = text_light

    # ==================== SLIDE 1: Title ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Title Card Shape
    shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = card_bg
    shape.line.color.rgb = accent_gold
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.6)

    p0 = tf.paragraphs[0]
    p0.text = "HACKATHON ACADÉMIQUE IPM 2026"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = accent_cyan
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "Système Intelligent de Transport Public\net de Paiement Numérique"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = text_light
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "GROUPE 3 : PAYMENT SERVICE (SPRING BOOT & ANGULAR)"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = accent_gold
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Gestion des Portefeuilles | Simulation Wave (+221 Sénégal) | Débit Trajet Bus | Détection de Fraude"
    p3.font.size = Pt(13)
    p3.font.color.rgb = text_muted
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(16)

    # ==================== SLIDE 2: Context & Mission ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "1. Mission du Groupe 3 & Contexte du Projet")

    # Left Box - Mission
    box1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = card_bg
    box1.line.color.rgb = accent_cyan

    tf1 = box1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.3)
    tf1.margin_top = Inches(0.3)

    p = tf1.paragraphs[0]
    p.text = "🎯 MISSION DU GROUPE 3"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_gold

    items1 = [
        "Développer le microservice responsable de la gestion des paiements et des soldes virtuels.",
        "Gérer les portefeuilles électroniques associés aux cartes de transport HAFILAT.",
        "Permettre la recharge rapide via la simulation Wave Mobile Money (+221 Sénégal).",
        "Traiter les demandes de débit envoyées par le Transport Service (Groupe 2)."
    ]
    for item in items1:
        p = tf1.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_before = Pt(10)

    # Right Box - Workflow
    box2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    box2.fill.solid()
    box2.fill.fore_color.rgb = card_bg
    box2.line.color.rgb = accent_green

    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)

    p = tf2.paragraphs[0]
    p.text = "🔄 RÔLE DANS L'ARCHITECTURE GLOBAL"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_green

    items2 = [
        "Étape 3 : L'usager recharge son portefeuille via Wave (+221).",
        "Étape 6 : Le passager valide sa carte/QR Code dans le bus (Transport Service).",
        "Étape 7 : Le Payment Service reçoit la demande et vérifie le solde disponible.",
        "Étape 8 : Débit automatique du tarif et enregistrement du reçu de transaction."
    ]
    for item in items2:
        p = tf2.add_paragraph()
        p.text = "✔ " + item
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_before = Pt(10)

    # ==================== SLIDE 3: Technical Features ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "2. Fonctionnalités Réalisées (Obligatoires & Bonus)")

    # 3 Cards Layout
    cards_data = [
        ("CORE FEATURES", accent_cyan, [
            "Recharge Portefeuille via Wave Mobile Money (Sénégal +221).",
            "Débit automatique de bus déclenché par NFC/QR Code.",
            "Vérification de solde & gestion d'erreur HTTP 402.",
            "Calculateur tarifaire automatique par ligne et profil."
        ]),
        ("OPTIONS BONUS INCLUSES", accent_gold, [
            "🚨 Détection de Fraude par IA/Règles (Passages trop rapprochés, tentatives répétées).",
            "📊 Tableau de bord analytique des revenus et transactions.",
            "💸 Gestion des remboursements (Refund API)."
        ]),
        ("STACK TECHNIQUE", accent_green, [
            "Backend : Java 17/21, Spring Boot 3, Spring Data JPA.",
            "Sécurité : Spring Security & JWT Filter.",
            "Frontend : Angular 17+ (Design Abu Dhabi Bus).",
            "Documentation : Swagger UI OpenAPI 3.0."
        ])
    ]

    for idx, (title, color, bullets) in enumerate(cards_data):
        x = Inches(0.8 + idx * 3.9)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        for b in bullets:
            p = tf.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(12)
            p.font.color.rgb = text_light
            p.space_before = Pt(10)

    # ==================== SLIDE 4: Architecture & API ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "3. Architecture API REST & Base de Données")

    # Table of REST Endpoints
    rows, cols = 6, 4
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8)
    table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(3.8)
    table.columns[2].width = Inches(4.4)
    table.columns[3].width = Inches(2.0)

    headers = ["Méthode", "Endpoint REST", "Description Fonctionnelle", "Statut / Port"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = card_bg
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = accent_gold

    api_data = [
        ("POST", "/api/payments/wave/recharge", "Recharge instantanée solde via simulation Wave (+221)", "201 Created"),
        ("POST", "/api/payments/process-ride", "Validation et débit automatique passager bus (NFC/QR)", "200 OK / 402"),
        ("POST", "/api/fares/calculate", "Calculateur automatique tarif selon ligne & catégorie", "200 OK"),
        ("GET", "/api/fraud/alerts", "Consultation des alertes de sécurité et fraudes détectées", "200 OK (Bonus)"),
        ("GET", "/api/stats/dashboard", "Tableau de bord des métriques financières globales", "200 OK (Bonus)")
    ]

    for row_idx, row_data in enumerate(api_data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_dark
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(12)
            p.font.color.rgb = text_light
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = accent_cyan

    # ==================== SLIDE 5: Demo & Deliverables ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "4. Démonstration Fonctionnelle & Livrables Fournis")

    box_demo = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    box_demo.fill.solid()
    box_demo.fill.fore_color.rgb = card_bg
    box_demo.line.color.rgb = accent_gold

    tf = box_demo.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.4)

    p = tf.paragraphs[0]
    p.text = "🎬 ÉTAPES DE LA DÉMONSTRATION EN DIRECT DEVENT LE JURY"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_gold

    steps = [
        "1. Lancement de l'interface Angular (http://localhost:4200) : Carte virtuelle HAFILAT & solde initial.",
        "2. Simulation Recharge Wave (+221 77 000 00 00) : Crédit instantané du portefeuille de 5 000 XOF.",
        "3. Simulation Trajet Bus (Interaction Groupe 2) : Validation du passager et débit automatique du solde.",
        "4. Gestion du Solde Insuffisant & Alerte Fraude : Blocage HTTP 402 et journalisation dans l'onglet Sécurité.",
        "5. Consultation Swagger UI (http://localhost:8083/swagger-ui.html) : Documentation OpenAPI interactive."
    ]

    for s in steps:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_before = Pt(12)

    # Save presentation
    output_path = r"c:\Users\LENOVO\OneDrive\Bureau\java ee\docs\Presentation_Payment_Service_Groupe3.pptx"
    prs.save(output_path)
    print(f"PowerPoint généré avec succès : {output_path}")

if __name__ == "__main__":
    create_presentation()
